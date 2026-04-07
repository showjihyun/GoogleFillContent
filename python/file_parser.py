"""
file_parser.py: 로컬 파일 파싱 (PDF, PPTX, XLSX, DOCX, HWP)

Drive에서 다운로드한 바이너리 파일을 텍스트로 변환한다.
Google 네이티브 포맷이 아닌 업로드 파일에 사용.
"""

import io
import logging
import os
import tempfile

log = logging.getLogger("fillcontent")


def parse_pdf(data: bytes) -> str | None:
    """PDF 바이너리에서 텍스트 추출."""
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            result = "\n\n--- PAGE BREAK ---\n\n".join(pages)
            return result if result.strip() else None
    except Exception as e:
        log.debug(f"    PDF 파싱 실패: {e}")
        return None


def parse_pptx(data: bytes) -> str | None:
    """PPTX 바이너리에서 텍스트 추출."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            texts.append(row_text)
            if texts:
                slides.append(f"--- 슬라이드 {i + 1} ---\n" + "\n".join(texts))
        result = "\n\n".join(slides)
        return result if result.strip() else None
    except Exception as e:
        log.debug(f"    PPTX 파싱 실패: {e}")
        return None


def parse_xlsx(data: bytes) -> str | None:
    """XLSX 바이너리에서 텍스트 추출."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"--- 시트: {ws.title} ---\n" + "\n".join(rows))
        wb.close()
        result = "\n\n".join(sheets)
        return result if result.strip() else None
    except Exception as e:
        log.debug(f"    XLSX 파싱 실패: {e}")
        return None


def parse_xls(data: bytes) -> str | None:
    """XLS (구형 Excel) 바이너리에서 텍스트 추출. olefile 기반 최소 추출."""
    try:
        import olefile
        if not olefile.isOleFile(io.BytesIO(data)):
            return None
        ole = olefile.OleFileIO(io.BytesIO(data))
        texts = []
        for stream in ole.listdir():
            try:
                raw = ole.openstream(stream).read()
                decoded = raw.decode("utf-8", errors="ignore")
                cleaned = "".join(c for c in decoded if c.isprintable() or c in "\n\t")
                if len(cleaned.strip()) > 20:
                    texts.append(cleaned.strip())
            except Exception:
                pass
        ole.close()
        result = "\n".join(texts)
        return result if result.strip() else None
    except Exception as e:
        log.debug(f"    XLS 파싱 실패: {e}")
        return None


def parse_docx(data: bytes) -> str | None:
    """DOCX 바이너리에서 텍스트 추출."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        # 테이블도 추출
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        result = "\n".join(paragraphs)
        return result if result.strip() else None
    except Exception as e:
        log.debug(f"    DOCX 파싱 실패: {e}")
        return None


def parse_doc(data: bytes) -> str | None:
    """DOC (구형 Word) 바이너리에서 텍스트 추출."""
    try:
        import olefile
        if not olefile.isOleFile(io.BytesIO(data)):
            return None
        ole = olefile.OleFileIO(io.BytesIO(data))
        # WordDocument 스트림에서 텍스트 추출
        if ole.exists("WordDocument"):
            raw = ole.openstream("WordDocument").read()
            # CP949/UTF-8 시도
            for enc in ["utf-8", "cp949", "euc-kr", "latin-1"]:
                try:
                    decoded = raw.decode(enc, errors="ignore")
                    cleaned = "".join(c for c in decoded if c.isprintable() or c in "\n\t ")
                    if len(cleaned.strip()) > 50:
                        ole.close()
                        return cleaned.strip()
                except Exception:
                    continue
        ole.close()
        return None
    except Exception as e:
        log.debug(f"    DOC 파싱 실패: {e}")
        return None


def parse_hwp(data: bytes) -> str | None:
    """HWP 바이너리에서 텍스트 추출 (OLE 스트림 기반)."""
    try:
        import olefile
        import zlib

        if not olefile.isOleFile(io.BytesIO(data)):
            return None

        ole = olefile.OleFileIO(io.BytesIO(data))
        texts = []

        # HWP의 본문은 BodyText/Section* 스트림에 저장
        for entry in ole.listdir():
            entry_path = "/".join(entry)
            if entry_path.startswith("BodyText/Section"):
                try:
                    raw = ole.openstream(entry).read()
                    # HWP는 보통 zlib 압축
                    try:
                        decompressed = zlib.decompress(raw, -15)
                    except Exception:
                        decompressed = raw

                    # 텍스트 추출 (UTF-16LE)
                    text_parts = []
                    i = 0
                    while i < len(decompressed):
                        # HWP 텍스트 레코드 파싱 (간소화)
                        if i + 1 < len(decompressed):
                            char = decompressed[i:i + 2]
                            code = int.from_bytes(char, "little")
                            # 제어 문자 건너뛰기
                            if code < 32 and code not in (10, 13, 9):
                                i += 2
                                continue
                            try:
                                text_parts.append(char.decode("utf-16-le", errors="ignore"))
                            except Exception:
                                pass
                        i += 2

                    section_text = "".join(text_parts)
                    # 제어 문자 정리
                    cleaned = "".join(
                        c for c in section_text
                        if c.isprintable() or c in "\n\t "
                    )
                    if cleaned.strip():
                        texts.append(cleaned.strip())
                except Exception:
                    pass

        ole.close()
        result = "\n".join(texts)
        return result if len(result.strip()) > 20 else None
    except Exception as e:
        log.debug(f"    HWP 파싱 실패: {e}")
        return None


# MIME 타입 → 파서 매핑
PARSERS = {
    "application/pdf": parse_pdf,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": parse_xlsx,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
    "application/msword": parse_doc,
    "application/vnd.ms-excel": parse_xls,
    "application/x-hwp": parse_hwp,
}


def can_parse(mime_type: str) -> bool:
    """로컬 파싱이 가능한 MIME 타입인지 확인."""
    return mime_type in PARSERS


def parse_file(data: bytes, mime_type: str) -> str | None:
    """MIME 타입에 맞는 파서로 바이너리 데이터를 텍스트로 변환."""
    parser = PARSERS.get(mime_type)
    if not parser:
        return None
    return parser(data)
